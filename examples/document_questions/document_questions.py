import argparse
import socketio
from pathlib import Path
from lollms import MSG_TYPE
from lollms.helpers import ASCIIColors
from lollms.paths import LollmsPaths
import time 
import json
from pathlib import Path
import PyPDF2
from docx import Document
from bs4 import BeautifulSoup
import json
import csv
from pptx import Presentation
import sys

# Connect to the Socket.IO server
sio = socketio.Client()


        
def read_pdf_file(file_path):
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def read_docx_file(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text


def read_json_file(file_path):
    with open(file_path, 'r') as file:
        data = json.load(file)
    return data


def read_csv_file(file_path):
    with open(file_path, 'r') as file:
        csv_reader = csv.reader(file)
        lines = [row for row in csv_reader]
    return lines    


def read_html_file(file_path):
    with open(file_path, 'r') as file:
        soup = BeautifulSoup(file, 'html.parser')
        text = soup.get_text()
    return text

def read_pptx_file(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

def read_text_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content

def chunk(input_text, word_count):
    # Check for empty text case (word count = 0)
    if not input_text or word_count == 0:
        return []
    
    # Split the string into words using .split() method. Remove leading and trailing whitespace characters from each word by replacing them with an empty strinng ""        
    cleaned_words = [word.strip(' ') for word in input_text.lower().split()] 
                     
    return [' '.join(cleaned_words[i:i+word_count]) for i in range(0,len(cleaned_words),word_count)]

# Event handler for receiving generated text
@sio.event
def text_generated(data):
    print('Generated text:', data)

def test_generate_text(host, port, lollms_paths:LollmsPaths):
    docs=lollms_paths.personal_data_path/"docs.txt"
    questions_path = lollms_paths.personal_data_path/"questions.txt"
    outputs=lollms_paths.personal_data_path/"outputs.txt"

    if not docs.exists():
        sys.exit(0)

    if not questions_path.exists():
        sys.exit(0)

    files = []
    # Read the text file and split by multiple newlines
    print("Loading files")
    with open(docs, 'r') as file:
        file = file.read().split('\n')
        files.append(Path(file[0]))

    for file in files:
        infos={
            "is_ready":False,
            "answer":""
        }
        if file.suffix.lower()==".pdf":
            txt = read_pdf_file(file)
        
        file_chunks = chunk(txt,512)
        # Event handler for successful connection
        @sio.event
        def connect():
            print('Connected to Socket.IO server')
            questions = []
            with open(questions_path, 'r') as file:
                questions = file.read().split('\n')
                
            for question in questions:
                print(f"Question:{question}")
                useful_chunks=[]
                for chunk in file_chunks:
                    prompt="Document:\n"+chunk+"\n"+question
                    # Trigger the 'generate_text' event with the prompt
                    infos["is_ready"]=False
                    print(f"Sending prompt:{prompt}")
                    sio.emit('generate_text', {'prompt': prompt, 'personality':-1, "n_predicts":1024})
                    while infos["is_ready"]==False:
                        time.sleep(0.1)
                    if(infos["answer"].lower()=="yes"):
                        ASCIIColors.info("Found useful chunk")
                        useful_chunks.append(chunk)
                with open(outputs, 'w') as file:
                    file.writelines([question,"Useful chunks"])
                    for i in range(len(useful_chunks)):
                        file.writelines([useful_chunks[i]])
                        

        @sio.event
        def text_chunk(data):
            print(data["chunk"],end="",flush=True)

        @sio.event
        def text_generated(data):
            print("text_generated_ok")
            print(data["text"])
            infos["answer"]=data["text"]
            infos["is_ready"]=True

        print(f"Connecting to http://{host}:{port}")
        # Connect to the Socket.IO server
        sio.connect(f'http://{host}:{port}')

        # Start the event loop
        sio.wait()

if __name__ == '__main__':
    # Parse command-line arguments
    parser = argparse.ArgumentParser(description='Socket.IO endpoint test')
    parser.add_argument('--host', type=str, default='localhost', help='Socket.IO server host')
    parser.add_argument('--port', type=int, default=9601, help='Socket.IO server port')
    parser.add_argument('--text-file', type=str, default=str(Path(__file__).parent/"example_text_gen.txt"),help='Path to the text file')
    args = parser.parse_args()
    lollms_paths = LollmsPaths.find_paths(force_local=False)
    # Run the test with provided arguments
    test_generate_text(args.host, args.port, lollms_paths)
